#!/usr/bin/env python3
"""
build_slides.py — Claude-themed, image-rich .pptx builder from slides_content.json.

Usage:
  python3 tools/build_slides.py --input temp/outputs/slides_content.json
  python3 tools/build_slides.py --input temp/outputs/slides_content.json --output path/to/deck.pptx
"""

import argparse
import hashlib
import io
import json
import os
import re
import ssl
import sys
import urllib.request

from dotenv import load_dotenv
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

load_dotenv()

TAG = "[build_slides]"

# ── Claude brand palette ──────────────────────────────────────────────────────
BG        = RGBColor(0x1A, 0x17, 0x14)   # warm dark charcoal
BG_PANEL  = RGBColor(0x26, 0x22, 0x1F)   # raised surface
ORANGE    = RGBColor(0xD9, 0x77, 0x06)   # amber — Claude primary
GOLD      = RGBColor(0xF5, 0x9E, 0x0B)   # golden highlight
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)   # violet secondary
TEAL      = RGBColor(0x0F, 0x96, 0x87)   # teal accent
CORAL     = RGBColor(0xE5, 0x60, 0x4A)   # coral / warm red
CREAM     = RGBColor(0xFB, 0xF4, 0xE3)   # primary text
WARM_GRAY = RGBColor(0xC8, 0xBA, 0xA6)   # secondary text
DIM       = RGBColor(0x78, 0x70, 0x6A)   # dim labels

FONT = os.getenv("SLIDE_FONT", "Calibri")
W = Inches(13.333)
H = Inches(7.5)

LEFT_FRAC  = 0.585
LEFT_W     = Inches(13.333 * LEFT_FRAC)
RIGHT_W    = Inches(13.333 * (1 - LEFT_FRAC))
RIGHT_LEFT = LEFT_W

# Accent palette cycles across content slides for visual variety
PALETTE = [ORANGE, GOLD, PURPLE, TEAL, CORAL]


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _seed(text):
    return int(hashlib.md5(text.encode()).hexdigest()[:8], 16) % 9999


def fetch_image(seed, img_w, img_h, cache_dir):
    fname = os.path.join(cache_dir, f"img_{seed}_{img_w}x{img_h}.jpg")
    if os.path.exists(fname):
        with open(fname, "rb") as f:
            return io.BytesIO(f.read())
    url = f"https://picsum.photos/seed/{seed}/{img_w}/{img_h}"
    try:
        ctx = ssl.create_default_context()
        ctx.check_hostname = False
        ctx.verify_mode = ssl.CERT_NONE
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=10, context=ctx) as resp:
            data = resp.read()
        with open(fname, "wb") as f:
            f.write(data)
        print(f"{TAG}   [img] downloaded seed={seed} → {os.path.basename(fname)}")
        return io.BytesIO(data)
    except Exception as e:
        print(f"{TAG}   [img] fetch failed (seed={seed}): {e}")
        return None


def _bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or BG


def _rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _text(slide, text, left, top, w, h, size,
          bold=False, italic=False, color=None, align=PP_ALIGN.LEFT):
    color = color or CREAM
    txb = slide.shapes.add_textbox(left, top, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _progress_bar(slide, current, total, accent):
    track_h = Pt(5)
    _rect(slide, 0, 0, W, track_h, BG_PANEL)
    if total > 0 and current > 0:
        fill_w = int(W * current / total)
        _rect(slide, 0, 0, fill_w, track_h, accent)


def _dot_cluster(slide, left, top, colors):
    for i, c in enumerate(colors):
        _rect(slide, left, top + i * Inches(0.22), Inches(0.14), Inches(0.14), c)


def _bottom_bar(slide, accent, width=None):
    w = width or LEFT_W + Pt(10)
    _rect(slide, 0, H - Pt(10), w, Pt(10), accent)
    _rect(slide, Inches(3.2), H - Pt(10), Inches(2.0), Pt(10), GOLD)
    _rect(slide, Inches(5.2), H - Pt(10), Inches(1.1), Pt(10), PURPLE)


# ── Slide builders ────────────────────────────────────────────────────────────

def build_title_slide(prs, slide_data, cache_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    title    = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle") or ""

    # Right image panel
    img = fetch_image(_seed(title), 560, 750, cache_dir)
    if img:
        slide.shapes.add_picture(img, RIGHT_LEFT, 0, RIGHT_W, H)
    else:
        _rect(slide, RIGHT_LEFT, 0, RIGHT_W, H, BG_PANEL)

    # Left dark panel (overlaps slightly to hide image edge)
    _rect(slide, 0, 0, LEFT_W + Pt(12), H, BG)

    # Left accent stripe
    _rect(slide, 0, 0, Pt(8), H, ORANGE)

    # Decorative dot cluster
    _dot_cluster(slide, Inches(0.46), Inches(0.48), [ORANGE, GOLD, PURPLE])

    # Label
    _text(slide, "CONFIDENTIAL  ·  INVESTMENT PRESENTATION",
          Inches(0.56), Inches(0.42), Inches(7.0), Inches(0.45),
          9.5, bold=True, color=ORANGE)

    # Title
    _text(slide, title,
          Inches(0.56), Inches(1.3), Inches(7.1), Inches(3.1),
          46, bold=True, color=CREAM)

    # Divider
    _rect(slide, Inches(0.56), Inches(4.3), Inches(4.2), Pt(3), ORANGE)

    # Subtitle
    if subtitle:
        _text(slide, subtitle,
              Inches(0.56), Inches(4.5), Inches(7.1), Inches(2.1),
              19, color=WARM_GRAY)

    _bottom_bar(slide, ORANGE)


def build_content_slide(prs, slide_data, accent, slide_num, total_slides,
                        content_idx, total_content, cache_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    title   = slide_data.get("title", "")
    bullets = slide_data.get("bullets", [])

    _progress_bar(slide, slide_num, total_slides, accent)

    # Right image panel
    _rect(slide, RIGHT_LEFT, Pt(5), RIGHT_W, H - Pt(5), BG_PANEL)
    img = fetch_image(_seed(title), 560, 730, cache_dir)
    if img:
        slide.shapes.add_picture(img, RIGHT_LEFT, Pt(5), RIGHT_W, H - Pt(5))

    # Left content panel
    _rect(slide, 0, Pt(5), LEFT_W, H - Pt(5), BG)
    _rect(slide, 0, Pt(5), Pt(6), H - Pt(5), accent)

    # Section number badge
    _rect(slide, Inches(0.46), Inches(0.28), Inches(0.56), Inches(0.5), accent)
    _text(slide, f"{content_idx:02d}",
          Inches(0.46), Inches(0.28), Inches(0.56), Inches(0.5),
          15, bold=True, color=BG, align=PP_ALIGN.CENTER)

    # Title
    _text(slide, title,
          Inches(1.16), Inches(0.22), Inches(6.45), Inches(0.82),
          28, bold=True, color=CREAM)

    # Rule
    _rect(slide, Inches(0.46), Inches(1.1), Inches(7.1), Pt(2), accent)

    # Bullets
    txb = slide.shapes.add_textbox(Inches(0.52), Inches(1.22),
                                   Inches(7.1), Inches(6.0))
    tf = txb.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(9)

        dot = p.add_run()
        dot.text = "● "
        dot.font.name = FONT
        dot.font.size = Pt(9)
        dot.font.color.rgb = accent

        body = p.add_run()
        body.text = bullet
        body.font.name = FONT
        body.font.size = Pt(19)
        body.font.color.rgb = CREAM

    # Slide counter
    _text(slide, f"{content_idx} / {total_content}",
          Inches(0.46), Inches(7.12), Inches(2), Inches(0.3),
          10, color=DIM)


def build_closing_slide(prs, slide_data, cache_dir):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    title    = slide_data.get("title", "Thank You")
    subtitle = slide_data.get("subtitle") or ""

    img = fetch_image(_seed(title + "_closing"), 560, 750, cache_dir)
    if img:
        slide.shapes.add_picture(img, RIGHT_LEFT, 0, RIGHT_W, H)
    else:
        _rect(slide, RIGHT_LEFT, 0, RIGHT_W, H, BG_PANEL)

    _rect(slide, 0, 0, LEFT_W + Pt(12), H, BG)
    _rect(slide, 0, 0, Pt(8), H, ORANGE)

    _dot_cluster(slide, Inches(0.46), Inches(0.48), [ORANGE, TEAL, GOLD])

    _text(slide, "NEXT STEPS",
          Inches(0.56), Inches(0.42), Inches(6), Inches(0.45),
          9.5, bold=True, color=ORANGE)

    _text(slide, title,
          Inches(0.56), Inches(1.3), Inches(6.9), Inches(2.8),
          42, bold=True, color=CREAM)

    _rect(slide, Inches(0.56), Inches(4.2), Inches(4.2), Pt(3), ORANGE)

    if subtitle:
        _text(slide, subtitle,
              Inches(0.56), Inches(4.4), Inches(6.9), Inches(2.4),
              17, color=WARM_GRAY)

    _bottom_bar(slide, ORANGE)


# ── Main ──────────────────────────────────────────────────────────────────────

def slugify(text):
    text = text.lower()
    text = re.sub(r"[^a-z0-9\s]", "", text)
    text = re.sub(r"\s+", "_", text.strip())
    return text[:40]


def main():
    parser = argparse.ArgumentParser(description="Build a Claude-themed .pptx")
    parser.add_argument("--input",  required=True, help="Path to slides_content.json")
    parser.add_argument("--output", default=None,  help="Output .pptx path (optional)")
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"{TAG} ERROR — input not found: {args.input}")
        sys.exit(1)

    print(f"{TAG} Loading slide content from: {args.input}")

    with open(args.input, encoding="utf-8") as f:
        try:
            data = json.load(f)
        except json.JSONDecodeError as e:
            print(f"{TAG} ERROR — bad JSON: {e}")
            sys.exit(1)

    deck_title = data.get("deck_title", "Presentation")
    accent_hex = data.get("accent_color", "#D97706")
    slides     = data.get("slides", [])

    if not re.match(r"^#[0-9A-Fa-f]{6}$", accent_hex):
        print(f"{TAG} ERROR — invalid accent_color: {accent_hex!r}")
        sys.exit(1)

    cache_dir = os.path.join(
        os.path.dirname(os.path.abspath(args.input)),
        "..", "resources", "images"
    )
    os.makedirs(cache_dir, exist_ok=True)

    total_slides  = len(slides)
    total_content = sum(1 for s in slides if s.get("type") == "content")
    content_idx   = 0

    print(f'{TAG} Deck title: "{deck_title}" | Slides: {total_slides} | Theme: Claude')
    print(f"{TAG} Image cache: {os.path.abspath(cache_dir)}")

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for i, sd in enumerate(slides):
        slide_type = sd.get("type")
        print(f"{TAG} Building slide {i + 1}/{total_slides} ({slide_type})...")

        if slide_type == "title":
            build_title_slide(prs, sd, cache_dir)

        elif slide_type == "content":
            content_idx += 1
            accent = PALETTE[(content_idx - 1) % len(PALETTE)]
            build_content_slide(prs, sd, accent, i + 1, total_slides,
                                content_idx, total_content, cache_dir)

        elif slide_type == "closing":
            build_closing_slide(prs, sd, cache_dir)

        else:
            print(f"{TAG} WARNING — unknown slide type '{slide_type}', skipping")

    out_dir  = os.path.dirname(os.path.abspath(args.input))
    out_path = args.output or os.path.join(out_dir, f"slides_{slugify(deck_title)}.pptx")
    prs.save(out_path)
    print(f"{TAG} SUCCESS — wrote {out_path}")


if __name__ == "__main__":
    main()
